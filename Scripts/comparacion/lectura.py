from pptx import Presentation
import os
from params import *

def read_ppt_content(ppt_path):
    try:
        ppt = Presentation(ppt_path)
    except Exception as e:
        print(f"Error al abrir el archivo PPT: {e}")
        return

    for slide_num, slide in enumerate(ppt.slides, start=1):
        print(f"\nDiapositiva número: {slide_num}\n{'='*20}")
        
        for shape in slide.shapes:
            # Texto en cuadros de texto
            if shape.has_text_frame:
                print("Cuadro de texto:")
                for paragraph in shape.text_frame.paragraphs:
                    print(paragraph.text)
            
            # Texto en tablas
            if shape.has_table:
                print("Tabla:")
                table = shape.table
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells]
                    print("\t".join(row_text))
            
            
            if shape.shape_type == 1:  # Placeholder (general type for SmartArt and others)
                if hasattr(shape, "text_frame") and shape.text_frame is not None:
                    print("SmartArt/Placeholder texto:")
                    for paragraph in shape.text_frame.paragraphs:
                        print(paragraph.text)

# Main

ppt_file_path = os.path.join(PPTS_PATH,"observaciones_modificado.pptx")
read_ppt_content(ppt_file_path)