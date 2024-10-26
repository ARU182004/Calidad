from pptx import Presentation
import os
import pandas as pd
from params import *

def extract_fields_from_ppt(ppt_path,ppt_filename):
    slides_data = []
    project_code = ""

    try:
        ppt = Presentation(ppt_path)
    except Exception as e:
        print(f"Error al abrir el archivo PPT: {e}")
        return slides_data

    for slide_num, slide in enumerate(ppt.slides, start=1):
        slide_data = {"slide_number": slide_num}

        if slide_num == 1:
            
            evaluacion = ""
            codigo_proyecto = ""
            text_boxes = [shape for shape in slide.shapes if shape.has_text_frame]

            if len(text_boxes) > 3:
                evaluacion = text_boxes[3].text_frame.text.strip()
            if len(text_boxes) > 5:
                codigo_proyecto = text_boxes[5].text_frame.text.strip()
                project_code = codigo_proyecto

            slide_data["evaluacion"] = evaluacion
            slide_data["codigo_proyecto"] = codigo_proyecto

        elif slide_num == 3:
            
            table_data_1 = []
            table_header_1 = ["Critico", "Alto", "Relevante", "Moderado", "Bajo", "Total"]
            specific_table_1_found = False

            
            table_data_2 = []
            table_header_2 = ["Proceso/Unidad auditada", "Calificativo", "Efectividad", "N° controles"]
            specific_table_2_found = False

            
            table_data_3 = []
            table_header_3 = ["Nro", "Título de la Observación", "Criticidad"]
            specific_table_3_found = False

            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                   
                    first_row_text = [cell.text for cell in table.rows[0].cells]
                    if first_row_text[0] == "Efectividad del control" and first_row_text[1] == "Criticidad del riesgo residual" and first_row_text[-1] == "Total":
                        specific_table_1_found = True
                        for row_idx in range(1, len(table.rows)):  
                            row_text = [table.cell(row_idx, col_idx).text for col_idx in range(len(table.columns))]
                            table_data_1.append(row_text)
                    elif first_row_text[:4] == table_header_2:
                        specific_table_2_found = True
                        for row_idx in range(1, len(table.rows)):  
                            row_text = [table.cell(row_idx, col_idx).text for col_idx in range(len(table.columns))]
                            table_data_2.append(row_text)
                    elif first_row_text == table_header_3:
                        specific_table_3_found = True
                        for row_idx in range(1, len(table.rows)):  
                            row_text = [table.cell(row_idx, col_idx).text for col_idx in range(len(table.columns))]
                            table_data_3.append(row_text)

            if specific_table_1_found:
                
                for row in table_data_1:
                    row_dict = {}
                    row_dict["efectividad"] = row[0]
                    for i, header in enumerate(table_header_1):
                        row_dict[f"{row[0].lower().replace(' ', '')}{header.lower()}"] = row[i + 1] if i + 1 < len(row) else 0
                    slide_data.update(row_dict)

            if specific_table_2_found:
                
                for idx, row in enumerate(table_data_2, start=1):
                    if idx > 3: 
                        break
                    slide_data[f"unidad_auditada_{idx}"] = row[0]
                    slide_data[f"calificativo_{idx}"] = row[1]
                    slide_data[f"efectividad_{idx}"] = row[2]
                    slide_data[f"numero_de_controles_{idx}"] = row[3]

            if specific_table_3_found:
                
                for idx, row in enumerate(table_data_3, start=1):
                    slide_data[f"observacion_{idx}"] = row[1]  
                    slide_data[f"criticidad_{idx}"] = row[2]  

        elif slide_num == 4:
            
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    
                    first_row_text = [cell.text for cell in table.rows[0].cells]
                    if "Campo" in first_row_text and "Dato" in first_row_text:
                        for row_idx in range(1, len(table.rows)):  
                            row = table.rows[row_idx]
                            hito = row.cells[0].text.strip()
                            fecha = row.cells[1].text.strip()
                            if hito == "Fecha de inicio":
                                slide_data["fecha_inicio"] = fecha
                            elif hito == "Fecha de término":
                                slide_data["fecha_termino"] = fecha
                            elif hito == "Fecha de emisión de Borrador de Informe":
                                slide_data["fecha_emision_borrador_informe"] = fecha
                            elif hito == "Fecha de emisión de Informe Final":
                                slide_data["fecha_emision_informe_final"] = fecha
                            elif hito == "Inicio Trabajo de campo":
                                slide_data["fecha_inicio_trabajo_campo"] = fecha
                            elif hito == "Fin Trabajo de campo":
                                slide_data["fecha_fin_trabajo_campo"] = fecha
                            elif hito == "IFC Total Inicial":
                                slide_data["IFC_total_inicial"] = fecha
                            elif hito == "IFC Total Residual":
                                slide_data["IFC_total_residual"] = fecha

        elif slide_num == 8:
            
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    first_row_text = [cell.text for cell in table.rows[0].cells]
                    
                    if "Equipo de Auditoría Interna" in first_row_text:
                        equipo_auditoria = {}
                        
                        for row_idx in range(1, len(table.rows)):  
                            row = table.rows[row_idx]
                            rol = row.cells[0].text.strip()
                            nombre = row.cells[1].text.strip()
                            if rol in equipo_auditoria:
                                equipo_auditoria[rol].append(nombre)
                            else:
                                equipo_auditoria[rol] = [nombre]
                        slide_data.update(equipo_auditoria)

        else:
            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = "\n".join([para.text for para in shape.text_frame.paragraphs])
                    print(f"Slide {slide_num} - Cuadro de texto: {text}")
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = [cell.text for cell in row.cells]
                        print(f"Slide {slide_num} - Tabla: {''.join(row_text)}")

        slides_data.append(slide_data)

    return slides_data, project_code

def export_to_excel(slides_data, project_code, ppt_filename, output_excel_path):
    records = []

    for slide in slides_data:
        slide_number = slide["slide_number"]
        for key, value in slide.items():
            if key != "slide_number":
                record = {
                    "Proyecto": project_code,
                    "Archivo": ppt_filename,
                    "Slide": slide_number,
                    "Campo encontrado": key,
                    "Valor de campo": value
                }
                records.append(record)

    df = pd.DataFrame(records)
    df.to_excel(output_excel_path, index=False)



ppt_file_path = os.path.join(PPTS_PATH, "informe_modificado.pptx")
ppt_filename = os.path.basename(ppt_file_path)
slides_data, project_code = extract_fields_from_ppt(ppt_file_path, ppt_filename)


for slide in slides_data:
    print(slide)


output_excel_path = os.path.join(RESULTS_PATH, "informe_result.xlsx")
export_to_excel(slides_data, project_code, ppt_filename, output_excel_path)

