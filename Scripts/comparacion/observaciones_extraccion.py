from pptx import Presentation
import os
from params import *
import pandas as pd

def extract_observations_from_ppt(ppt_path, ppt_filename):
    observations = []
    project_code = ""

    try:
        ppt = Presentation(ppt_path)
    except Exception as e:
        print(f"Error al abrir el archivo PPT: {e}")
        return observations

    for slide_num, slide in enumerate(ppt.slides, start=1):
        slide_data = {"slide_number": slide_num}
        previous_row_text = []

        if slide_num == 1:
            
            evaluacion = ""
            codigo_proyecto = ""
            text_boxes = [shape for shape in slide.shapes if shape.has_text_frame]

            if len(text_boxes) > 3:
                evaluacion = text_boxes[3].text_frame.text.strip()
            if len(text_boxes) > 5:
                codigo_proyecto = text_boxes[5].text_frame.text.strip()
                project_code = codigo_proyecto

            slide_data["evaluacion"] = evaluacion
            slide_data["codigo_proyecto"] = codigo_proyecto
        
        else:

            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    num_rows = len(table.rows)
                    num_cols = len(table.columns)

                   
                    if num_rows > 1:
                        for row in table.rows:
                            row_text = [cell.text.strip() for cell in row.cells]

                            if "Hallazgo Control" in row_text[0]:
                                slide_data[f"hallazgo_control_{slide_num}"] = row_text[0]
                            elif "Título de la Observación" in row_text[0]:
                                slide_data[f"titulo_observacion_{slide_num}"] = ""
                                previous_row_text = ["titulo_observacion", slide_num]
                            elif "Observación" in row_text[0]:
                                slide_data[f"observacion_{slide_num}"] = ""
                                previous_row_text = ["observacion", slide_num]
                            elif "Recomendación" in row_text[0]:
                                slide_data[f"recomendacion_{slide_num}"] = ""
                                previous_row_text = ["recomendacion", slide_num]
                            elif "Plan de acción" in row_text[0]:
                                slide_data[f"plan_accion_{slide_num}"] = ""
                                previous_row_text = ["plan_accion", slide_num]
                            elif "Responsable" in row_text[0]:
                                print(row_text)
                                slide_data[f"responsable_{slide_num}"] = row_text[1]
                            elif "Cargo" in row_text[0]:
                                slide_data[f"cargo_{slide_num}"] = row_text[1]
                            elif "Estado" in row_text[0]:
                                slide_data[f"estado_{slide_num}"] = row_text[1]
                            elif "Fecha de vencimiento" in row_text[0]:
                                slide_data[f"fecha_vencimiento_{slide_num}"] = row_text[1]
                            elif previous_row_text:
                                field_name, slide_idx = previous_row_text
                                if field_name == 'titulo_observacion':
                                    slide_data[f"{field_name}_{slide_idx}"] = row_text[1]
                                else:
                                    slide_data[f"{field_name}_{slide_idx}"] = row_text[0]
                                previous_row_text = []

        observations.append(slide_data)

    return observations, project_code

def export_observations_to_excel(project_code,observations, output_excel_path):
    records = []

    for slide in observations:
        slide_number = slide["slide_number"]
        for key, value in slide.items():
            if key != "slide_number":
                record = {
                    "Proyecto": project_code,
                    "Archivo":ppt_filename,
                    "Slide": slide_number,
                    "Campo encontrado": key,
                    "Valor de campo": value
                }
                records.append(record)

    df = pd.DataFrame(records)
    df.to_excel(output_excel_path, index=False)




ppt_file_path = os.path.join(PPTS_PATH, "observaciones_modificado.pptx")
ppt_filename = os.path.basename(ppt_file_path)
observations,project_code = extract_observations_from_ppt(ppt_file_path, ppt_filename)
print(observations)


output_excel_path = os.path.join(RESULTS_PATH, "observaciones_extracted.xlsx")
export_observations_to_excel(project_code,observations, output_excel_path)
