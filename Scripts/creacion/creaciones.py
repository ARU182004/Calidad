from pptx import Presentation
import os
from pptx.util import Pt
from params import *
from functions import obtenerpath

def create_memorando(tipo,negocio,codigo,df, output_path):

    

    if negocio == 'SEGUROS':
        ppt_template_path = os.path.join(PLANTILLAS_PS_PATH, "07. Sprint Planning.pptx")
    elif negocio == 'PRIMA':
        ppt_template_path = os.path.join(PLANTILLAS_PRI_PATH, "07. Sprint Planning.pptx")
    
    elif negocio == 'SALUD':
        ppt_template_path = os.path.join(PLANTILLAS_EPS_PATH, "07. Sprint Planning.pptx")
    
    elif negocio == 'CREDISEGUROS':
        ppt_template_path = os.path.join(PLANTILLAS_CRE_PATH, "07. Sprint Planning.pptx")
    else:
        ppt_template_path = os.path.join(PLANTILLAS_PS_PATH, "07. Sprint Planning.pptx")
    


    try:
        ppt = Presentation(ppt_template_path)
    except Exception as e:
        print(f"Error al abrir el archivo PPT: {e}")
        return

    for slide_num, slide_data in enumerate(memorando_list, start=1):
        if slide_num > len(ppt.slides):
            print(f"No hay suficiente cantidad de slides en la presentación para el diccionario número {slide_num}")
            break
        slide = ppt.slides[slide_num - 1]
        print(f"Diapositiva número: {slide_num}")

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, column_name in slide_data.items():
                            if run.text == key:
                                run.text = str(df[column_name].iloc[0])

    try:
        final_path = obtenerpath(tipo,negocio,codigo)
        ppt.save(final_path)
        print(f"PPT guardado en {final_path}")
    except Exception as e:
        print(f"Error al guardar el archivo PPT: {e}")




def create_observaciones_ppt(tipo, negocio, codigo, df, output_path):
    num_slides = len(df) 
    
    if negocio == 'SEGUROS':
        ppt_template_path = os.path.join(PLANTILLAS_PS_PATH, f"plantilla_observaciones_{num_slides}.pptx")
    elif negocio == 'PRIMA':
        ppt_template_path = os.path.join(PLANTILLAS_PRI_PATH, f"plantilla_observaciones_{num_slides}.pptx")
    
    elif negocio == 'SALUD':
        ppt_template_path = os.path.join(PLANTILLAS_EPS_PATH, f"plantilla_observaciones_{num_slides}.pptx")
    
    elif negocio == 'CREDISEGUROS':
        ppt_template_path = os.path.join(PLANTILLAS_CRE_PATH, f"plantilla_observaciones_{num_slides}.pptx")
    else:
        ppt_template_path = os.path.join(PLANTILLAS_PS_PATH, f"plantilla_observaciones_{num_slides}.pptx")

    try:
        ppt = Presentation(ppt_template_path)
    except Exception as e:
        print(f"Error al abrir el archivo PPT: {e}")
        return

   
    first_slide = ppt.slides[0]
    for shape in first_slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for key, value in primera_slide_dict.items():
                        if key in run.text:
                            run.text = run.text.replace(key, str(value))

    
    font_size_observaciones = Pt(14)

    
    for index, row in enumerate(df.itertuples(index=False)):
        slide = ppt.slides[index + 1]  
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, column_name in observaciones_list[0].items():
                            if key in run.text:
                                if key == 'Unidad_auditada_n':
                                    run.text = run.text.replace(key, str(getattr(row, column_name)))
                                else:
                                    run.text = run.text.replace(key, str(getattr(row, column_name)))
                                    run.font.size = font_size_observaciones

        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for cell in table.iter_cells():
                    for key, column_name in observaciones_list[0].items():
                        if key in cell.text:
                            cell.text = cell.text.replace(key, str(getattr(row, column_name)))
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.size = font_size_observaciones

    try:
        final_path = obtenerpath(tipo, negocio, codigo)
        ppt.save(final_path)
        print(f"PPT guardado en {final_path}")
    except Exception as e:
        print(f"Error al guardar el archivo PPT: {e}")


def create_informe_ppt(tipo,negocio,codigo,df_datos, df_calificativo_total, df_calificativo_unidad_responsable, df_cantidad_controles, df_observaciones_informe, output_path):
   

    if negocio == 'SEGUROS':
        ppt_template_path = os.path.join(PLANTILLAS_PS_PATH, "01. Borrador Informe.pptx")
    elif negocio == 'PRIMA':
        ppt_template_path = os.path.join(PLANTILLAS_PRI_PATH, "01. Borrador Informe.pptx")
    
    elif negocio == 'SALUD':
        ppt_template_path = os.path.join(PLANTILLAS_EPS_PATH, "01. Borrador Informe.pptx")
    
    elif negocio == 'CREDISEGUROS':
        ppt_template_path = os.path.join(PLANTILLAS_CRE_PATH, "01. Borrador Informe.pptx")
    else:
        ppt_template_path = os.path.join(PLANTILLAS_PS_PATH, "01. Borrador Informe.pptx")

    try:
        ppt = Presentation(ppt_template_path)
    except Exception as e:
        print(f"Error al abrir el archivo PPT: {e}")
        return

    
    listas_diccionarios = [
        (calificativo_total_list, df_calificativo_total),
        (calificativo_unidad_responsable_list, df_calificativo_unidad_responsable),
        (cantidad_controles_list, df_cantidad_controles),
        (observaciones_informe_list, df_observaciones_informe),
        (memorando_list, df_datos)
    ]

    
    font_size_first_slide = Pt(40)  
    font_size_other_slides = Pt(18)  

    
    for lista_dicc, df in listas_diccionarios:
        for slide_num, slide_data in enumerate(lista_dicc):
            if slide_num >= len(ppt.slides):
                break
            slide = ppt.slides[slide_num]

            
            font_size = font_size_first_slide if slide_num == 0 else font_size_other_slides

            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            for key, column_name in slide_data.items():
                                if lista_dicc in [calificativo_unidad_responsable_list, observaciones_informe_list]:
                                    max_records = 3 if lista_dicc == observaciones_informe_list else 2
                                    for i in range(min(len(df), max_records)):  # Solo tomar en cuenta los primeros n registros
                                        new_key = f"{key}_{i + 1}"
                                        if new_key in run.text:
                                            run.text = run.text.replace(new_key, str(df[column_name].iloc[i]))
                                            run.font.size = font_size
                                else:
                                    if key in run.text:
                                        run.text = run.text.replace(key, str(df[column_name].iloc[0]))
                                        run.font.size = font_size

            
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    for cell in table.iter_cells():
                        for key, column_name in slide_data.items():
                            if lista_dicc in [calificativo_unidad_responsable_list, observaciones_informe_list]:
                                max_records = 3 if lista_dicc == observaciones_informe_list else 2
                                for i in range(min(len(df), max_records)):  # Solo tomar en cuenta los primeros n registros
                                    new_key = f"{key}_{i + 1}"
                                    if new_key in cell.text:
                                        cell.text = cell.text.replace(new_key, str(df[column_name].iloc[i]))
                                        for paragraph in cell.text_frame.paragraphs:
                                            for run in paragraph.runs:
                                                run.font.size = font_size
                            else:
                                if key in cell.text:
                                    cell.text = cell.text.replace(key, str(df[column_name].iloc[0]))
                                    for paragraph in cell.text_frame.paragraphs:
                                        for run in paragraph.runs:
                                            run.font.size = font_size

    try:
        final_path = obtenerpath(tipo,negocio,codigo)
        ppt.save(final_path)
        print(f"PPT guardado en {final_path}")
    except Exception as e:
        print(f"Error al guardar el archivo PPT: {e}")