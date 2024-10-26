from pptx import Presentation
from params import *
import os
import re

phrases_to_remove = [
    "Gerencia Área de Auditoría Interna",
    "Pacífico Seguros",
    "Evaluación Programada",
    "Junio 2024"
]

def clean_text(text):
    for phrase in phrases_to_remove:
        text = text.replace(phrase, "")
    return text.replace(":", "").strip()

def extract_specific_content(ppt_path):
    try:
        ppt = Presentation(ppt_path)
    except Exception as e:
        print(f"Error al abrir el archivo PPT: {e}")
        return

    ppt_content = {}

    project_code_pattern = re.compile(r'\b[A-Z]{1,3}-[A-Z]{2,4}-\d{2,3}-\d{4}\b')

    for slide_num, slide in enumerate(ppt.slides, start=1):
        slide_content = {}

        if slide_num == 1:
            slide_content["evaluation_name"] = ""
            slide_content["project_code"] = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    text = text.replace('\x0b', ' ').replace('\t', ' ')
                    lines = text.split('\n')
                    for line in lines:
                        line = clean_text(line.strip())
                        if line:
                            # Buscando el código del proyecto
                            match = project_code_pattern.search(line)
                            if match:
                                slide_content["project_code"] = match.group(0)
                                # Eliminar el código del proyecto del nombre de la evaluación
                                line = line.replace(match.group(0), "").strip()
                            # Asignar el resto del texto como el nombre de la evaluación
                            if line:
                                slide_content["evaluation_name"] = line
            ppt_content[slide_num] = slide_content
        
        elif slide_num == 3:
            slide_content["gerencias_calificadas"] = []
            slide_content["evaluacion_controles"] = []
            slide_content["principales_observaciones"] = []
            slide_content["rangos"] = []

            tables = []
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    table_content = []
                    for row in table.rows:
                        row_text = [cell.text for cell in row.cells]
                        table_content.append(row_text)
                    tables.append(table_content)
            

            # Asignamos las tablas a las llaves correspondientes
            if len(tables) >= 3:
                slide_content["gerencias_calificadas"] = tables[2]
                slide_content["evaluacion_controles"] = tables[3]
                slide_content["principales_observaciones"] = tables[1]
                
                

            ppt_content[slide_num] = slide_content

        elif slide_num == 4:
            slide_content["fecha_inicio"] = ""
            slide_content["fecha_termino"] = ""
            slide_content["fecha_emision_borrador"] = ""
            slide_content["trabajo_campo"] = ""
            slide_content["ifc_total_inicial"] = ""
            slide_content["ifc_total_residual"] = ""
            slide_content["normas_iai"] = ""

            normas_iai_started = False
            normas_iai_text = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    lines = text.split('\n')
                    for line in lines:
                        line = clean_text(line.strip())
                        
                    
                        if "Normas IAI" in line:
                            normas_iai_started = True
                        if normas_iai_started and "Cuadro de texto:" not in line:
                            normas_iai_text.append(line)
                elif shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            
                            text = cell.text
                            
                            lines = text.split('\n')
                            for line in lines:
                                line = clean_text(line.strip())
                                print(line)
                                print("-----")
                                if "Fecha de inicio" in line:
                                    fecha_inicio_match = re.search(r'Fecha de inicio (\d{2}/\d{2}/\d{4})', line)
                                    fecha_termino_match = re.search(r'Fecha de término (\d{2}/\d{2}/\d{4})', line)
                                    
                                    if fecha_inicio_match:
                                        slide_content["fecha_inicio"] = fecha_inicio_match.group(1)
                                    if fecha_termino_match:
                                        slide_content["fecha_termino"] = fecha_termino_match.group(1)
                                if "Fecha de emisión de Borrador de Informe" in line:
                                    fecha_emision_borrador_match = re.search(r'Fecha de emisión de Borrador de Informe (\d{2}\.\d{2}\.\d{4})', line)
                                    if fecha_emision_borrador_match:
                                        slide_content["fecha_emision_borrador"] = fecha_emision_borrador_match.group(1)
                                if "Trabajo de campo" in line:
                                    trabajo_campo_match = re.search(r'Trabajo de campo(.+)', line)
                                    ifc_total_inicial_match = re.search(r'IFC Total Inicial\s+(\d+)', line)
                                    ifc_total_residual_match = re.search(r'IFC Total Residual\s+(\d+)', line)
                                    if trabajo_campo_match:
                                        slide_content["trabajo_campo"] = trabajo_campo_match.group(1)
                                    if ifc_total_inicial_match:
                                        slide_content["ifc_total_inicial"] = ifc_total_inicial_match.group(1)
                                    if ifc_total_residual_match:
                                        slide_content["ifc_total_residual"] = ifc_total_residual_match.group(1)
                                if "Normas IAI" in line:
                                    normas_iai_started = True
                                if normas_iai_started and "Cuadro de texto:" not in line:
                                    normas_iai_text.append(line)

            slide_content["normas_iai"] = " ".join(normas_iai_text).strip()
            ppt_content[slide_num] = slide_content

    return ppt_content

# Main
ppt_file_path = os.path.join(PPTS_PATH, "informe_borrador.pptx")

ppt_content = extract_specific_content(ppt_file_path)

# Imprimir el contenido extraído para ver el resultado
for slide_num, content in ppt_content.items():
    print(f"\nDiapositiva número: {slide_num}")
    if slide_num == 1:
        print(f"Nombre de la Evaluación: {content['evaluation_name']}")
        print(f"Código del Proyecto: {content['project_code']}")
    elif slide_num == 3:
        print("Gerencias Calificadas y del Proceso:")
        for row in content["gerencias_calificadas"]:
            print("\t".join(row))
        print("Evaluación de Controles:")
        for row in content["evaluacion_controles"]:
            print("\t".join(row))
        print("Principales Observaciones:")
        for row in content["principales_observaciones"]:
            print("\t".join(row))
    elif slide_num == 4:
        print(f"Fecha de inicio: {content['fecha_inicio']}")
        print(f"Fecha de término: {content['fecha_termino']}")
        print(f"Fecha de emisión de Borrador de Informe: {content['fecha_emision_borrador']}")
        print(f"Trabajo de campo: {content['trabajo_campo']}")
        print(f"IFC Total Inicial: {content['ifc_total_inicial']}")
        print(f"IFC Total Residual: {content['ifc_total_residual']}")
        print(f"Normas IAI: {content['normas_iai']}")